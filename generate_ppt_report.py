from __future__ import annotations

import argparse
import logging
import math
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from generate_report import (
    FundNavData,
    compute_metrics,
    load_config,
    load_fund_nav,
    pct,
    report_datetime,
    risk_label,
    setup_logging,
    stable_seed,
    synthetic_nav_series,
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG = RGBColor(248, 246, 239)
INK = RGBColor(26, 32, 38)
MUTED = RGBColor(101, 112, 125)
LINE = RGBColor(218, 213, 201)
BLUE = RGBColor(38, 99, 179)
TEAL = RGBColor(25, 137, 116)
RED = RGBColor(191, 64, 64)
AMBER = RGBColor(191, 132, 45)
PAPER = RGBColor(255, 253, 248)


class SlideBuilder:
    def __init__(self, prs: Presentation, title: str, kicker: str | None = None) -> None:
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.slide.background.fill.solid()
        self.slide.background.fill.fore_color.rgb = BG
        if kicker:
            self.text(0.55, 0.33, 2.2, 0.25, kicker.upper(), size=8, color=MUTED, bold=True)
        self.text(0.55, 0.58, 8.8, 0.45, title, size=24, color=INK, bold=True)
        self.line(0.55, 1.14, 12.25, 0.0, LINE)

    def text(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        *,
        size: int = 12,
        color: RGBColor = INK,
        bold: bool = False,
        align: int | None = None,
    ):
        box = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        p = frame.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = value
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def rect(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: RGBColor = PAPER,
        line: RGBColor | None = LINE,
    ):
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(0.5)
        return shape

    def line(self, x: float, y: float, w: float, h: float, color: RGBColor = LINE):
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(max(w, 0.01)), Inches(max(h, 0.01)))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape


def collect_rows(config: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for fund in config["funds"]:
        try:
            nav_data = load_fund_nav(fund)
            metrics = compute_metrics(nav_data.nav)
        except Exception as exc:
            logging.exception("PPT metric generation failed for %s, using synthetic fallback: %s", fund.get("name"), exc)
            fallback_nav = synthetic_nav_series(seed=stable_seed(str(fund["id"])))
            nav_data = FundNavData(
                nav=fallback_nav,
                source="synthetic_nav_series",
                latest_date=fallback_nav.index[-1].strftime("%Y-%m-%d"),
            )
            metrics = compute_metrics(nav_data.nav)

        rows.append(
            {
                "id": fund["id"],
                "name": fund["name"],
                "code": str(fund.get("fund_code", "")),
                "category": fund.get("category", ""),
                "risk_level": fund.get("risk_level", ""),
                "nav": float(nav_data.nav.iloc[-1]),
                "nav_date": nav_data.latest_date,
                "daily": float(metrics.daily_return),
                "weekly": float(metrics.weekly_return),
                "monthly": float(metrics.monthly_return),
                "drawdown": float(metrics.max_drawdown),
                "volatility": float(metrics.volatility),
                "source": nav_data.source,
                "status": risk_label(metrics),
            }
        )
    return rows


def color_for_return(value: float) -> RGBColor:
    if pd.isna(value):
        return MUTED
    return TEAL if value >= 0 else RED


def color_for_status(status: str) -> RGBColor:
    if status == "高风险关注":
        return RED
    if status == "中等波动":
        return AMBER
    return TEAL


def draw_metric_card(builder: SlideBuilder, x: float, y: float, w: float, h: float, label: str, value: str, note: str, color: RGBColor) -> None:
    builder.rect(x, y, w, h, PAPER, LINE)
    builder.text(x + 0.22, y + 0.18, w - 0.44, 0.22, label, size=9, color=MUTED, bold=True)
    builder.text(x + 0.22, y + 0.47, w - 0.44, 0.42, value, size=24, color=color, bold=True)
    builder.text(x + 0.22, y + 0.98, w - 0.44, 0.28, note, size=9, color=MUTED)


def draw_return_bars(builder: SlideBuilder, rows: list[dict[str, Any]], key: str, x: float, y: float, title: str) -> None:
    builder.text(x, y, 3.7, 0.24, title, size=11, color=INK, bold=True)
    max_abs = max(max(abs(row[key]) for row in rows), 0.01)
    center = x + 1.95
    bar_max = 1.45
    for index, row in enumerate(rows):
        yy = y + 0.42 + index * 0.42
        builder.text(x, yy - 0.02, 1.18, 0.22, row["name"][:8], size=8, color=MUTED)
        builder.line(center, yy + 0.06, 0.01, 0.13, LINE)
        width = max(abs(row[key]) / max_abs * bar_max, 0.02)
        if row[key] >= 0:
            builder.rect(center, yy + 0.04, width, 0.16, color_for_return(row[key]), None)
            label_x = center + width + 0.05
        else:
            builder.rect(center - width, yy + 0.04, width, 0.16, color_for_return(row[key]), None)
            label_x = center + 0.05
        builder.text(label_x, yy - 0.01, 0.55, 0.22, pct(row[key]), size=8, color=INK)


def add_cover(prs: Presentation, rows: list[dict[str, Any]], report_date: str) -> None:
    slide = SlideBuilder(prs, f"基金组合日报 {report_date}", "fund tracker")
    high_risk = sum(1 for row in rows if row["status"] == "高风险关注")
    fallback = sum(1 for row in rows if row["source"] == "synthetic_nav_series")
    avg_daily = sum(row["daily"] for row in rows) / len(rows)
    avg_monthly = sum(row["monthly"] for row in rows) / len(rows)

    slide.text(0.68, 1.55, 5.8, 0.55, "真实净值驱动的每日组合快照", size=28, color=INK, bold=True)
    slide.text(0.7, 2.15, 6.6, 0.42, "含收益表现、回撤波动、数据源状态与高风险基金提示。", size=14, color=MUTED)
    draw_metric_card(slide, 0.7, 3.05, 2.55, 1.45, "平均日涨跌", pct(avg_daily), "五只基金简单平均", color_for_return(avg_daily))
    draw_metric_card(slide, 3.55, 3.05, 2.55, 1.45, "平均近1月", pct(avg_monthly), "五只基金简单平均", color_for_return(avg_monthly))
    draw_metric_card(slide, 6.4, 3.05, 2.55, 1.45, "高风险关注", str(high_risk), "回撤或波动触发", RED if high_risk else TEAL)
    draw_metric_card(slide, 9.25, 3.05, 2.55, 1.45, "Fallback", str(fallback), "真实数据失败数量", RED if fallback else TEAL)

    slide.text(0.7, 5.15, 2.8, 0.28, "今日风险判断", size=11, color=INK, bold=True)
    if high_risk:
        note = "存在高风险关注基金，建议复核 QDII 汇率、海外科技股估值与短期回撤。"
        note_color = RED
    else:
        note = "当前未触发高风险状态，继续执行每日跟踪。"
        note_color = TEAL
    slide.rect(0.7, 5.52, 11.1, 0.8, PAPER, LINE)
    slide.text(0.95, 5.74, 10.45, 0.34, note, size=16, color=note_color, bold=True)


def add_return_slide(prs: Presentation, rows: list[dict[str, Any]]) -> None:
    slide = SlideBuilder(prs, "收益表现：近1月分化明显，QDII 贡献主要波动", "returns")
    draw_return_bars(slide, rows, "daily", 0.7, 1.5, "当日涨跌幅")
    draw_return_bars(slide, rows, "weekly", 4.65, 1.5, "近1周收益")
    draw_return_bars(slide, rows, "monthly", 8.6, 1.5, "近1个月收益")

    slide.text(0.7, 4.55, 2.2, 0.24, "基金净值状态", size=11, color=INK, bold=True)
    headers = ["基金", "最新净值", "净值日期", "数据源"]
    col_x = [0.7, 4.95, 6.65, 8.55]
    widths = [4.0, 1.45, 1.55, 1.55]
    for x, w, header in zip(col_x, widths, headers):
        slide.text(x, 4.9, w, 0.2, header, size=8, color=MUTED, bold=True)
    for index, row in enumerate(rows):
        y = 5.25 + index * 0.33
        slide.text(col_x[0], y, widths[0], 0.2, f"{row['name']} ({row['code']})", size=8, color=INK)
        slide.text(col_x[1], y, widths[1], 0.2, f"{row['nav']:.4f}", size=8, color=INK)
        slide.text(col_x[2], y, widths[2], 0.2, row["nav_date"], size=8, color=INK)
        slide.text(col_x[3], y, widths[3], 0.2, row["source"], size=8, color=TEAL if row["source"] == "Eastmoney" else RED)


def add_risk_slide(prs: Presentation, rows: list[dict[str, Any]]) -> None:
    slide = SlideBuilder(prs, "风险图谱：回撤与年化波动共同决定关注优先级", "risk map")
    slide.text(0.7, 1.42, 5.0, 0.25, "最大回撤", size=11, color=INK, bold=True)
    slide.text(7.0, 1.42, 5.0, 0.25, "年化波动率", size=11, color=INK, bold=True)
    max_drawdown = max(abs(row["drawdown"]) for row in rows) or 0.01
    max_vol = max(row["volatility"] for row in rows) or 0.01
    for index, row in enumerate(rows):
        y = 1.9 + index * 0.65
        status_color = color_for_status(row["status"])
        slide.text(0.7, y - 0.02, 2.25, 0.24, row["name"][:10], size=9, color=INK)
        draw_w = max(abs(row["drawdown"]) / max_drawdown * 2.6, 0.05)
        slide.rect(3.05, y, draw_w, 0.22, RED, None)
        slide.text(5.78, y - 0.02, 0.75, 0.24, pct(row["drawdown"]), size=8, color=INK)

        slide.text(7.0, y - 0.02, 2.25, 0.24, row["name"][:10], size=9, color=INK)
        vol_w = max(row["volatility"] / max_vol * 2.6, 0.05)
        slide.rect(9.35, y, vol_w, 0.22, status_color, None)
        slide.text(12.08, y - 0.02, 0.75, 0.24, pct(row["volatility"]), size=8, color=INK)

    slide.rect(0.7, 5.65, 11.6, 0.7, PAPER, LINE)
    high_risk_names = "、".join(row["name"] for row in rows if row["status"] == "高风险关注") or "无"
    slide.text(0.95, 5.84, 10.9, 0.28, f"高风险关注：{high_risk_names}", size=14, color=RED if high_risk_names != "无" else TEAL, bold=True)


def add_source_slide(prs: Presentation, rows: list[dict[str, Any]], report_date: str) -> None:
    slide = SlideBuilder(prs, "数据源与下载说明", "delivery")
    fallback = [row for row in rows if row["source"] == "synthetic_nav_series"]
    source_note = "五只基金均使用 Eastmoney/Tiantian Fund 真实净值。" if not fallback else "部分基金触发 fallback，请检查日志。"
    slide.text(0.75, 1.55, 10.8, 0.4, source_note, size=20, color=TEAL if not fallback else RED, bold=True)
    slide.text(0.75, 2.15, 10.8, 0.35, f"PPTX 文件路径：reports/{report_date}.pptx", size=14, color=INK)
    slide.text(0.75, 2.58, 10.8, 0.35, f"Markdown 文件路径：reports/{report_date}.md", size=14, color=INK)

    slide.text(0.75, 3.35, 2.3, 0.25, "数据新鲜度", size=11, color=INK, bold=True)
    for index, row in enumerate(rows):
        y = 3.8 + index * 0.42
        slide.text(0.75, y, 3.0, 0.22, row["name"], size=9, color=INK)
        slide.text(4.05, y, 1.5, 0.22, row["nav_date"], size=9, color=MUTED)
        slide.text(5.8, y, 1.6, 0.22, row["source"], size=9, color=TEAL if row["source"] == "Eastmoney" else RED)
        slide.text(7.55, y, 2.2, 0.22, row["status"], size=9, color=color_for_status(row["status"]), bold=True)

    slide.text(0.75, 6.28, 10.8, 0.25, "提示：本报告仅用于投资记录和风险监测，不构成收益保证或买卖建议。", size=9, color=MUTED)


def build_ppt(config: dict[str, Any], output_dir: Path) -> Path:
    report_date = report_datetime(config).strftime("%Y-%m-%d")
    rows = collect_rows(config)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover(prs, rows, report_date)
    add_return_slide(prs, rows)
    add_risk_slide(prs, rows)
    add_source_slide(prs, rows, report_date)

    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{report_date}.pptx"
    prs.save(output_path)
    return output_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate daily fund PPTX report")
    parser.add_argument("--config", default="config.yaml")
    parser.add_argument("--output-dir", default="reports")
    args = parser.parse_args()

    setup_logging(Path("logs"))
    config = load_config(Path(args.config))
    output_path = build_ppt(config, Path(args.output_dir))
    print(f"Generated PPTX report: {output_path}")


if __name__ == "__main__":
    main()
